from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE
from pptx.oxml.shapes.shared import qn
import json
import os

class JSONToPowerPointConverter:
    def __init__(self):
        self.prs = None
        
    def create_presentation_from_json(self, json_file_path, template_path=None):
        """Create PowerPoint presentation from JSON data"""
        
        # Load JSON data
        with open(json_file_path, 'r', encoding='utf-8') as f:
            presentation_data = json.load(f)
        
        # Create presentation (use template if provided)
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            # Clear existing slides
            slide_count = len(self.prs.slides)
            for i in range(slide_count - 1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        else:
            self.prs = Presentation()
        
        # Set presentation dimensions if available
        if 'slide_width' in presentation_data:
            self.prs.slide_width = presentation_data['slide_width']
        if 'slide_height' in presentation_data:
            self.prs.slide_height = presentation_data['slide_height']
        
        # Create slides from JSON data
        for slide_data in presentation_data.get('slides', []):
            self._create_slide_from_data(slide_data)
        
        return self.prs
    
    def _create_slide_from_data(self, slide_data):
        """Create a slide from JSON slide data"""
        
        # Get appropriate layout (use blank layout for maximum control)
        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)
        
        # Apply background if specified
        if 'background' in slide_data and slide_data['background']:
            self._apply_slide_background(slide, slide_data['background'])
        
        # Add shapes to slide
        for shape_data in slide_data.get('shapes', []):
            self._create_shape_from_data(slide, shape_data)
    
    def _apply_slide_background(self, slide, background_data):
        """Apply background to slide"""
        try:
            if background_data.get('type') == 'solid' and background_data.get('color'):
                color_info = background_data['color']
                if color_info and color_info.get('type') == 'rgb':
                    rgb = color_info['rgb']
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        except Exception as e:
            print(f"Warning: Could not apply background: {e}")
    
    def _create_shape_from_data(self, slide, shape_data):
        """Create a shape from JSON shape data"""
        
        shape_type = shape_data.get('shape_type')
        position = shape_data.get('position', {})
        
        # Extract position and size
        left = position.get('left', 0)
        top = position.get('top', 0)
        width = position.get('width', Inches(2))
        height = position.get('height', Inches(1))
        
        created_shape = None
        
        # Create shape based on type
        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape_data.get('text_content'):
            created_shape = self._create_textbox_from_data(slide, shape_data, left, top, width, height)
        
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            created_shape = self._create_autoshape_from_data(slide, shape_data, left, top, width, height)
        
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            created_shape = self._create_picture_from_data(slide, shape_data, left, top, width, height)
        
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            created_shape = self._create_table_from_data(slide, shape_data, left, top, width, height)
        
        # Apply common shape properties
        if created_shape:
            self._apply_shape_properties(created_shape, shape_data)
    
    def _create_textbox_from_data(self, slide, shape_data, left, top, width, height):
        """Create text box from shape data"""
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Apply text content and formatting
        text_content = shape_data.get('text_content')
        if text_content:
            self._apply_text_formatting(text_frame, text_content)
        
        return textbox
    
    def _apply_text_formatting(self, text_frame, text_content):
        """Apply text formatting from JSON data"""
        
        # Clear existing text
        text_frame.clear()
        
        # Apply text frame properties
        if 'margin_left' in text_content:
            text_frame.margin_left = text_content['margin_left']
        if 'margin_top' in text_content:
            text_frame.margin_top = text_content['margin_top']
        if 'margin_right' in text_content:
            text_frame.margin_right = text_content['margin_right']
        if 'margin_bottom' in text_content:
            text_frame.margin_bottom = text_content['margin_bottom']
        if 'word_wrap' in text_content:
            text_frame.word_wrap = text_content['word_wrap']
        if 'auto_size' in text_content:
            text_frame.auto_size = text_content['auto_size']
        
        # Add paragraphs
        paragraphs_data = text_content.get('paragraphs', [])
        
        if not paragraphs_data and text_content.get('text'):
            # Simple text without detailed formatting
            text_frame.text = text_content['text']
            return
        
        for i, para_data in enumerate(paragraphs_data):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            # Apply paragraph properties
            if 'alignment' in para_data and para_data['alignment'] is not None:
                paragraph.alignment = para_data['alignment']
            if 'level' in para_data:
                paragraph.level = para_data['level']
            
            # Add runs
            runs_data = para_data.get('runs', [])
            if not runs_data and para_data.get('text'):
                # Simple paragraph text
                paragraph.text = para_data['text']
                continue
            
            paragraph.clear()
            for run_data in runs_data:
                run = paragraph.add_run()
                run.text = run_data.get('text', '')
                
                # Apply font formatting
                font_data = run_data.get('font', {})
                self._apply_font_formatting(run.font, font_data)
    
    def _apply_font_formatting(self, font, font_data):
        """Apply font formatting from JSON data"""
        
        if font_data.get('name'):
            font.name = font_data['name']
        
        if font_data.get('size'):
            font.size = Pt(font_data['size'])
        
        if font_data.get('bold') is not None:
            font.bold = font_data['bold']
        
        if font_data.get('italic') is not None:
            font.italic = font_data['italic']
        
        if font_data.get('underline') is not None:
            font.underline = font_data['underline']
        
        # Apply color
        color_data = font_data.get('color')
        if color_data and color_data.get('type') == 'rgb':
            rgb = color_data['rgb']
            font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    
    def _create_autoshape_from_data(self, slide, shape_data, left, top, width, height):
        """Create auto shape from shape data"""
        
        # Create rectangle as default auto shape (you can expand this)
        autoshape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        
        # Apply text if present
        text_content = shape_data.get('text_content')
        if text_content:
            self._apply_text_formatting(autoshape.text_frame, text_content)
        
        return autoshape
    
    def _create_table_from_data(self, slide, shape_data, left, top, width, height):
        """Create table from shape data"""
        
        table_data = shape_data.get('table_data', {})
        rows = table_data.get('rows', 2)
        columns = table_data.get('columns', 2)
        
        # Create table
        table_shape = slide.shapes.add_table(rows, columns, left, top, width, height)
        table = table_shape.table
        
        # Populate table data
        data = table_data.get('data', [])
        for row_idx, row_data in enumerate(data):
            if row_idx >= rows:
                break
            for col_idx, cell_data in enumerate(row_data):
                if col_idx >= columns:
                    break
                
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_data.get('text', '')
                
                # Apply cell formatting
                if cell_data.get('text_formatting'):
                    self._apply_text_formatting(cell.text_frame, cell_data['text_formatting'])
        
        return table_shape
    
    def _create_picture_from_data(self, slide, shape_data, left, top, width, height):
        """Create picture placeholder (requires actual image file)"""
        
        # Note: This creates a placeholder rectangle since we don't have the actual image
        # In a real implementation, you'd need to store and restore image data
        
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Gray placeholder
        
        # Add text indicating it's an image placeholder
        text_frame = placeholder.text_frame
        text_frame.text = f"Image: {shape_data.get('name', 'Unknown')}"
        
        return placeholder
    
    def _apply_shape_properties(self, shape, shape_data):
        """Apply general shape properties"""
        
        # Apply rotation
        if 'rotation' in shape_data:
            shape.rotation = shape_data['rotation']
        
        # Apply fill
        fill_data = shape_data.get('fill')
        if fill_data and fill_data.get('solid_color'):
            color_info = fill_data['solid_color']
            if color_info and color_info.get('type') == 'rgb':
                rgb = color_info['rgb']
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        # Apply line
        line_data = shape_data.get('line')
        if line_data and hasattr(shape, 'line'):
            if line_data.get('color') and line_data['color'].get('type') == 'rgb':
                rgb = line_data['color']['rgb']
                shape.line.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            
            if line_data.get('width'):
                shape.line.width = line_data['width']
    
    def save_presentation(self, output_path):
        """Save the created presentation"""
        if self.prs:
            self.prs.save(output_path)
            print(f"Presentation created and saved to: {output_path}")
        else:
            print("No presentation to save. Create presentation first.")

# Utility function to create JSON from existing PowerPoint
def export_ppt_to_json(ppt_path, json_path):
    """Export existing PowerPoint to JSON format"""
    
    from artifacts import PowerPointEditor  # Assuming you have the previous class
    
    editor = PowerPointEditor(ppt_path)
    presentation_data = editor.read_complete_presentation()
    
    # Custom JSON encoder to handle special objects
    class CustomJSONEncoder(json.JSONEncoder):
        def default(self, obj):
            if hasattr(obj, '__dict__'):
                return obj.__dict__
            return str(obj)
    
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, indent=2, ensure_ascii=False, cls=CustomJSONEncoder)
    
    print(f"PowerPoint exported to JSON: {json_path}")

# Example usage and testing
def example_usage():
    """Example of how to use the JSON to PowerPoint converter"""
    
    # Create sample JSON data
    sample_json_data = {
        "slide_width": 9144000,  # Standard slide width in EMUs
        "slide_height": 6858000,  # Standard slide height in EMUs
        "slides": [
            {
                "slide_number": 1,
                "layout_name": "Title Slide",
                "background": {
                    "type": "solid",
                    "color": {"type": "rgb", "rgb": [240, 248, 255]}
                },
                "shapes": [
                    {
                        "shape_index": 0,
                        "shape_type": 17,  # TEXT_BOX
                        "name": "Title",
                        "position": {
                            "left": 914400,    # 1 inch in EMUs
                            "top": 1828800,    # 2 inches in EMUs
                            "width": 7315200,  # 8 inches in EMUs
                            "height": 1371600  # 1.5 inches in EMUs
                        },
                        "text_content": {
                            "text": "Sample Presentation",
                            "paragraphs": [
                                {
                                    "paragraph_index": 0,
                                    "text": "Sample Presentation",
                                    "alignment": 1,  # Center alignment
                                    "runs": [
                                        {
                                            "run_index": 0,
                                            "text": "Sample Presentation",
                                            "font": {
                                                "name": "Arial",
                                                "size": 44,
                                                "bold": True,
                                                "italic": False,
                                                "color": {"type": "rgb", "rgb": [0, 0, 139]}
                                            }
                                        }
                                    ]
                                }
                            ]
                        }
                    },
                    {
                        "shape_index": 1,
                        "shape_type": 17,  # TEXT_BOX
                        "name": "Subtitle",
                        "position": {
                            "left": 914400,
                            "top": 3657600,    # 4 inches
                            "width": 7315200,
                            "height": 914400   # 1 inch
                        },
                        "text_content": {
                            "text": "Created from JSON data",
                            "paragraphs": [
                                {
                                    "paragraph_index": 0,
                                    "text": "Created from JSON data",
                                    "alignment": 1,
                                    "runs": [
                                        {
                                            "run_index": 0,
                                            "text": "Created from JSON data",
                                            "font": {
                                                "name": "Arial",
                                                "size": 24,
                                                "bold": False,
                                                "italic": True,
                                                "color": {"type": "rgb", "rgb": [70, 130, 180]}
                                            }
                                        }
                                    ]
                                }
                            ]
                        }
                    }
                ]
            },
            {
                "slide_number": 2,
                "layout_name": "Content Slide",
                "shapes": [
                    {
                        "shape_index": 0,
                        "shape_type": 17,
                        "name": "Content",
                        "position": {
                            "left": 914400,
                            "top": 914400,
                            "width": 7315200,
                            "height": 4572000
                        },
                        "text_content": {
                            "text": "• Bullet point 1\n• Bullet point 2\n• Bullet point 3",
                            "paragraphs": [
                                {
                                    "paragraph_index": 0,
                                    "text": "• Bullet point 1",
                                    "runs": [
                                        {
                                            "run_index": 0,
                                            "text": "• Bullet point 1",
                                            "font": {
                                                "name": "Calibri",
                                                "size": 18,
                                                "bold": False,
                                                "color": {"type": "rgb", "rgb": [0, 0, 0]}
                                            }
                                        }
                                    ]
                                }
                            ]
                        }
                    }
                ]
            }
        ]
    }
    
    # Save sample JSON
    with open('sample_presentation.json', 'w') as f:
        json.dump(sample_json_data, f, indent=2)
    
    # Convert JSON to PowerPoint
    converter = JSONToPowerPointConverter()
    converter.create_presentation_from_json('sample_presentation.json')
    converter.save_presentation('from_json.pptx')

if __name__ == "__main__":
    example_usage()
    
    # Example of round-trip: PPT -> JSON -> PPT
    # Uncomment the following lines if you have an existing PowerPoint file
    
    # # Step 1: Export existing PowerPoint to JSON
    # export_ppt_to_json('existing_presentation.pptx', 'exported.json')
    # 
    # # Step 2: Create new PowerPoint from JSON
    # converter = JSONToPowerPointConverter()
    # converter.create_presentation_from_json('exported.json')
    # converter.save_presentation('recreated_from_json.pptx')
